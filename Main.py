from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_FILL
import aspose.slides as slides
import os
from os import listdir
from os.path import isfile, join

def main():
    print("type s to save old ppt's and modify in new folder")
    print("type m to modify ppt's")

    # mode = "s"
    mode=input()

    print("exemple:C:\\Users\\Rafael\\Desktop\\cantari ")
    print("source folder: ")
    # mypath=r"C:\Users\Rafael\Desktop\cantari"
    mypath = input()
    newpath=""
    if(mode == "s"):
        print("destination folder: ")
        newpath=input()
        # newpath=r"C:\Users\Rafael\Desktop\cantari2"

    else:
        if(mode == "m"):
            newpath=mypath

    onlyfiles = [f for f in listdir(mypath) if isfile(join(mypath, f))]
    # print(onlyfiles)
    for name in onlyfiles:
        if "ppt" not in name or "lnk" in name:
            print("excluding .. "+name)
            onlyfiles.remove(name)

    originalfiles=onlyfiles
    onlyfiles2=[]
    failedfiles=[]

    for file in onlyfiles:
        name = mypath + "\\" + file
        #conversie la pptx
        if "pptx" not in name:
            with slides.Presentation(name) as presentation:
                name2 = mypath + "\\" + file + "x"
                presentation.save(name2, slides.export.SaveFormat.PPTX)
                onlyfiles2.append(name2)
        name=name2
        file=file+"x"
        try:
            prs = Presentation(name)
            print("Opened "+file)
            for slide_layout in prs.slide_master.slide_layouts:
                fill = slide_layout.background.fill
                fill.solid()
                #fundal negru
                fill.fore_color.rgb = RGBColor(0, 0, 0)
            #setare 16:9
            prs.slide_width = Inches(16)
            prs.slide_height=Inches(9)
            for slide in prs.slides:
                for shape in slide.shapes:

                    # size = 18
                    # for paragraph in shape.text_frame.paragraphs:
                    #     for run in paragraph.runs:
                    #         if run.font.size:
                    #             size2=int(str(run.font.size))/12_700
                    #             if size<size2:
                    #                 size=size2

                    # print(shape.text_frame.text)
                    shape.left, shape.top = Inches(0.5), Inches(0.5)
                    #shape.fill.type = MSO_FILL.SOLID
                    shape.fill.transparency=1
                    # gaseste strofa
                    if(len(shape.text_frame.text)>10):
                        frame=shape.text_frame
                        #sterge textul in plus creat de conversia de la ppt la pptx
                        if("Evaluation" in frame.text):
                            frame.text=""
                        shape.width = Inches(15)
                        shape.height = Inches(6)
                        frame.fit_text('Times New Roman', max_size=90)

                        frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                        frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        #centreaza textul
                        shape.left, shape.top = int((prs.slide_width - shape.width) / 2),int((prs.slide_height - shape.height) / 2)
                    else:
                        # a gasit numarul slide-ului sau numarul cantarii
                        if("/" not in shape.text_frame.text): #daca nu e numarul slide-ului
                            if("efren" not in shape.text_frame.text): #daca nu e refren
                                shape.left, shape.top = Inches(14.5), Inches(8.5)

                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        #scris alb
                        paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            name = newpath + "\\" + file
            try:
                os.remove(name)
            except:
               pass
            prs.save(name)
            print("success..."+file)
        except:
            print("fail..."+file)
            failedfiles.append(file)

    #sterge elementele in plus
    if(mode == "m"):
        for file in originalfiles:
            os.remove(mypath + "\\" + file)
    else:
        for file in onlyfiles2:
            os.remove(file)

    print(failedfiles)

if __name__ == '__main__':
    main()
